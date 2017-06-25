import os
import sys
import csv
from collections import Counter, defaultdict

import pandas as pd
from statsmodels.stats.inter_rater import aggregate_raters, fleiss_kappa
from pptx import Presentation

# configure Django so we can use models from the annotate app 
sys.path.append('/home/nejl/Dropbox/projects/tator/repo/tator')
os.environ['DJANGO_SETTINGS_MODULE'] = 'tagit.settings'
import django
django.setup()

from django.contrib.auth.models import User

from annotate.models import Query, Annotation, UserResponse
from templates import slide_template


# TODO: need to add sampling method that samples equally from dividing the
# probability mass of the distribution into thirds: top most frequent, middle,
# and bottom.

# TODO: fix analysis to have a global collection filter and then make sure adding
# annotations to queries in a different collection does not change the results 

def split_data_frame_by_prob(df, column, nbins):
    # splits a dataframe into 'nbins' of equal probability mass
    # using column specified by 'coilumn'

    df = df.sort_values(column, ascending=False)
    values = df[column]
    
    bin_probability = 1/nbins
    total = sum(values)
    cutoffs = []
    cumulative_total = 0 
    next_bin_probability = bin_probability
    
    for i, count in enumerate(values):
        cumulative_total += count

        if cumulative_total/total < next_bin_probability:
            continue

        cutoffs.append(i)
        next_bin_probability += bin_probability

    start = 0
    new_dfs = []
    while cutoffs:
        cutoff = cutoffs.pop(0)
        if len(cutoffs) == 0:
            # last item; get the rest
            new_dfs.append(df[start:])
        else:
            new_dfs.append(df[start:cutoff])
        start = cutoff

    return new_dfs


def load_queries(path):
    with open(path) as csvfile:
        df = pd.read_csv(csvfile, delimiter=';')
    return df
    

def clean_queries(df):
    """Returns the input DataFrame of queries cleaned"""
    # filter out queries with length less than 2 characters long
    df = df[df['querystring'].str.len() > 1]
    return df


def split_num(num, splits):
    """Returns the number 'num' divided into a list of numbers of size 'splits' """
    splits = [int(num/splits)+1]*(num%splits) + [int(num/splits)]*(splits-num%splits)
    assert sum(splits) == num
    return splits


def import_queries(path, collection, sample='first', limit=None, allow_dupes=False): 
    df = load_queries(path)
    df = clean_queries(df)
    
    if not allow_dupes:
        # remove existing queries from candidate queries to sample
        existing = [query.text for query in Query.objects.all()]
        df = df[~df['querystring'].isin(existing)]

    if limit is not None:
        if sample == 'first':
            df = df[:limit]
        elif sample == 'random':
            df = df.sample(limit)
        elif sample == 'proportional':
            df = df.sample(limit, weights='countqstring')
        elif sample == 'split':
            split_size = 3
            splits = split_data_frame_by_prob(df, 'countqstring', split_size)
            sizes = split_num(limit, split_size)
            sub_samples = []
            for size, split_df in zip(sizes, splits): 
                sub_samples.append(split_df.sample(size, weights='countqstring'))
            df = pd.concat(sub_samples)
            assert len(df) == limit
        else:
            print('Unknown sampling method')
            return

    for i, values in enumerate(df.values.tolist()):
         text, count = values
         Query.objects.create(text=text, count=count, collection=collection)

    print("Added {} queries to the database.\n".format(i+1))
    print(df.describe())


def pretty_print_counter(counter, reverse=False):
    lines = []
    for key, value in sorted(counter.items(), reverse=reverse):
        lines.append("{}: {}".format(key, value))
    return "\n".join(lines)


def get_user_results(username):
    # for each user, display the number of results
    # user
    lines = ["*** Annotator: {} ***".format(username)]
    lines.append("===================================\n")
    responses = UserResponse.objects.filter(user__username=username)
    annotations = [r for r in responses if r.annotation]
    skipped = [r for r in responses if r.skipped]

    lines.append("{} Skipped Queries:\n".format(len(skipped)))
    for response in skipped:
        line ='    "{}"\n    --- "{}"'.format(response.query.text,
                                              response.skipped.description)
        lines.append(line)

    lines.append("\n{} Annotations:\n".format(len(annotations)))

    lines.append(Annotation._meta.get_field('is_geo').verbose_name)
    q1 = Counter(r.annotation.is_geo for r in annotations)
    lines.append(pretty_print_counter(q1, reverse=True))

    lines.append("")

    lines.append(Annotation._meta.get_field('loc_type').verbose_name)
    q2 = Counter(r.annotation.loc_type for r in annotations)
    lines.append(pretty_print_counter(q2, reverse=True))

    lines.append("")
    
    lines.append(Annotation._meta.get_field('query_type').verbose_name)
    q3 = Counter(r.annotation.query_type for r in annotations)
    lines.append(pretty_print_counter(q3))
    
    return "\n".join(lines)
    

def do_iaa_pairs(user_pairs, questions=(1,2,3), level='fine'):
    results = defaultdict(list)
    for question in questions:
        for users in user_pairs:
            kappa = get_iaa(question, users=users, level=level)
            results[question].append(kappa)
    return results


def print_iaa_pairs(results, user_pairs):
    print('    '+'   '.join(', '.join(user) for user in user_pairs))
    for question, kappas in results.items():
        ks = ''.join("{:0<5.3}                    ".format(k) for k in kappas)
        print("Q{}: {}".format(question, ks))
    

def get_iaa(question_num, queries=None, users=None, level='fine'):
    data  = get_annotations(question_num, queries, users, level)
    #n_cat = Annotation.get_num_categories(question_num)
    results = aggregate_raters(data, n_cat=None)
    kappa = fleiss_kappa(results[0])
    return kappa


def get_annotations(question_num, queries=None, users=None, level='fine'):
    assert level in ('fine', 'coarse')
        
    queries = Query.objects.exclude(responses__skipped__isnull=False).distinct()

    if queries is not None:
        queries = queries.filter(pk__in=queries)
        
    data = []
    for query in queries:
        # get all non-skipped results
        responses = query.responses.exclude(skipped__isnull=False)
        
        if users is not None:
            # restrict annotations to supplied users
            responses = responses.filter(user__username__in=users)

        results = [r.annotation.get_question(question_num) for r in responses]

        if question_num in (2,3) and level == 'coarse':
            # use course grained agreement
            results = [r[0] for r in results]
            
        if results:
            data.append(results)
    return data


def show_agreement(question_num, users, skip_agree=True):
    lines = []
    queries = Query.objects.exclude(responses__skipped__isnull=False).distinct()
    queries = sorted(queries, key=lambda x:x.pk)
    users.sort()
    col_width = max(len(u) for u in users) + 2
    
    lines.append("".join("{u:{width}}".format(u=u, width=col_width)
                         for u in users))
    agree = 0
    disagree = 0
    
    for query in queries:
        responses  = query.responses.order_by('user__username')
        answers = [r.annotation.get_question(question_num) for r in responses]
        if skip_agree and len(set(answers)) <= 1:
            # all annotators agree, skip
            agree += 1
            continue
        disagree += 1
        line = "".join("{a:<{width}}".format(a=a, width=col_width)
                       for a in answers) + query.text
        lines.append(line)

    start = [
        "Question {}:".format(question_num),
        "Number all agree: {}".format(agree),
        "Number with some disagreement: {}".format(disagree),
        ""
    ]

    return "\n".join(start + lines)


def get_results(users):
    queries = Query.objects.exclude(responses__skipped__isnull=False).distinct()
    queries = sorted(queries, key=lambda x:x.pk)
    users.sort()
    
    rest_cols = ["Q{}_{}".format(num, user) for user in users for num in (1,2,3)]
    header = ['id', 'query'] + rest_cols
    rows = [header]
    
    for query in queries:
        row = [query.pk, query.text]
        responses = query.responses.order_by('user__username')
        for response in responses:
            row.append(response.annotation.get_question(1))
            row.append(response.annotation.get_question(2))
            row.append(response.annotation.get_question(3))
        rows.append(row)
    return rows


def export_results_csv(users, outfile='annotations.csv'):
    results = get_results(users)
    with open(outfile, 'w', encoding='utf8', newline='') as csvfile:
        writer = csv.writer(csvfile, delimiter=',')
        writer.writerows(results)

        
def make_slides_latex(users, csv=None, outfile='slides/slides.tex'):
    if csv is None:
        results = get_results(users)
    else:
        with open(csv, encoding='utf8') as csvfile:
            results = list(csv.reader(csvfile))

    lines = []
    header = results[0]    
    for i, query in enumerate(results[1:]):
        row1 = r"Q1 & {} & {} & {}\\".format(query[2], query[5], query[8])
        row2 = r"Q2 & {} & {} & {}\\".format(query[3], query[6], query[9])
        row3 = r"Q3 & {} & {} & {}\\".format(query[4], query[7], query[10])
        rows = "\n".join([row1, row2, row3])
        title = "Query {}".format(i+1)
        slide = slide_template.format(title=title, query=query[1], rows=rows)
        lines.append(slide)
        
    with open(outfile, 'w', encoding='utf8') as texfile:
        texfile.write('\n'.join(lines))


def make_slides_pptx(users, csv=None):
    """Not finished. Used latex instead"""
    if csv is None:
        results = get_results(users)
    else:
        with open(csv, encoding='utf8') as csvfile:
            results = list(csv.reader(csvfile))
            
    header = results[0]

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
        
    for i, query in enumerate(results[1:]):
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = 'Query {}'.format(i+1)

        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame

        p = tf.paragraphs[0]
        p.text = query[1]
        p.level = 0
        

    prs.save('test.pptx')
